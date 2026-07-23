"""
生成"心理异常智能早筛与精准干预系统"项目汇报 PPT
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── 色彩体系 ──────────────────────────────────────────
DARK_BG     = RGBColor(0x1B, 0x3A, 0x4B)   # 深蓝绿
TEAL        = RGBColor(0x0D, 0x7C, 0x85)   # 主色
TEAL_LIGHT  = RGBColor(0x14, 0xA3, 0xA8)   # 浅青
MINT        = RGBColor(0x5C, 0xD6, 0xC3)   # 薄荷高亮
LIGHT_BG    = RGBColor(0xF0, 0xF5, 0xF4)   # 浅灰绿底
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT    = RGBColor(0x1B, 0x3A, 0x4B)
MUTED_TEXT   = RGBColor(0x6B, 0x82, 0x90)
CARD_BG     = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_ORANGE = RGBColor(0xE8, 0x6A, 0x3A)  # 危机/警告
ACCENT_GOLD  = RGBColor(0xD4, 0xA0, 0x1F)   # 路由
ROUTE_COMFORT = RGBColor(0x14, 0xA3, 0xA8)  # 安慰路
ROUTE_KNOWLEDGE = RGBColor(0x0D, 0x7C, 0x85) # 知识路
ROUTE_CRISIS = RGBColor(0xE8, 0x6A, 0x3A)   # 危机路

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── 工具函数 ──────────────────────────────────────────
def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text="",
                font_size=14, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei", margin_left=0, margin_right=0,
                margin_top=0, margin_bottom=0, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    # margins
    tf.margin_left = Inches(margin_left)
    tf.margin_right = Inches(margin_right)
    tf.margin_top = Inches(margin_top)
    tf.margin_bottom = Inches(margin_bottom)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    p.line_spacing = Pt(int(font_size * line_spacing))
    return txBox

def add_rich_textbox(slide, left, top, width, height, segments,
                     align=PP_ALIGN.LEFT, line_spacing=1.3, margin_left=0.1,
                     margin_right=0.1, margin_top=0.1, margin_bottom=0.1):
    """segments: list of dicts [{text,size,color,bold,font_name,break_line}]"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin_left)
    tf.margin_right = Inches(margin_right)
    tf.margin_top = Inches(margin_top)
    tf.margin_bottom = Inches(margin_bottom)

    for i, seg in enumerate(segments):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = seg.get("text", "")
        p.font.size = Pt(seg.get("size", 14))
        p.font.color.rgb = seg.get("color", DARK_TEXT)
        p.font.bold = seg.get("bold", False)
        p.font.name = seg.get("font_name", "Microsoft YaHei")
        p.alignment = align
        p.line_spacing = Pt(int(seg.get("size", 14) * line_spacing))
        if seg.get("break_line"):
            p = tf.add_paragraph()
    return txBox

def add_rect(slide, left, top, width, height, fill_color=WHITE,
             border_color=None, border_width=0, shadow=False, radius=None):
    if radius:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.adjustments[0] = radius
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    if shadow:
        shape.shadow.inherit = False
    return shape

def add_arrow_right(slide, left, top, width, height, color=TEAL):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_chevron(slide, left, top, width, height, color=TEAL):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def set_shape_text(shape, text, font_size=12, color=WHITE, bold=True,
                   font_name="Microsoft YaHei", align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align

def add_section_number(slide, left, top, number):
    """圆形编号"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(0.4), Inches(0.4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()
    set_shape_text(shape, str(number), font_size=14, color=WHITE, bold=True)

def add_page_title(slide, number, title):
    """统一的内容页标题"""
    add_section_number(slide, 0.6, 0.35, number)
    add_textbox(slide, 1.15, 0.3, 10, 0.5, title,
                font_size=28, color=DARK_TEXT, bold=True)

# ═════════════════════════════════════════════════════════
# SLIDE 1 — 封面
# ═════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BG)

# 顶部装饰条
add_rect(slide, 0, 0, 13.333, 0.06, MINT)

# 左侧大装饰圆
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(3.5), Inches(4), Inches(4))
shape.fill.solid()
shape.fill.fore_color.rgb = TEAL
shape.line.fill.background()
shape.fill.fore_color.brightness = 0.0  # opaque

# 主标题
add_textbox(slide, 2.5, 1.8, 9, 1.0,
            "心理异常智能早筛与精准干预系统",
            font_size=42, color=WHITE, bold=True, font_name="Microsoft YaHei")

# 副标题
add_textbox(slide, 2.5, 2.9, 9, 0.6,
            "基于多模态AI的心理咨询辅助系统",
            font_size=22, color=MINT, font_name="Microsoft YaHei")

# 分隔线
add_rect(slide, 2.5, 3.7, 2.0, 0.04, MINT)

# 细节信息
add_textbox(slide, 2.5, 4.0, 9, 0.5,
            "安全过滤  ·  情绪分析  ·  智能路由  ·  精准干预",
            font_size=16, color=RGBColor(0x9B, 0xB8, 0xBF), font_name="Microsoft YaHei")

add_textbox(slide, 2.5, 4.6, 9, 0.4,
            "FastAPI + Qwen2.5-7B + SenseVoice + Chroma RAG + MySQL",
            font_size=13, color=RGBColor(0x7A, 0x9E, 0xA8), font_name="Consolas")

# ═════════════════════════════════════════════════════════
# SLIDE 2 — 项目概述
# ═════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_page_title(slide, 1, "项目概述")

# 一句话定位
add_rect(slide, 0.6, 1.0, 12.1, 0.7, TEAL)
add_textbox(slide, 0.8, 1.05, 11.7, 0.55,
            "一个融合大语言模型（LLM）、语音识别（ASR）、面部情绪识别与知识检索（RAG）的 AI 心理治疗师系统",
            font_size=16, color=WHITE, bold=False)

# 四大核心能力卡片 - 2x2 网格
cards = [
    ("安全过滤", "S1", "关键词匹配 + 敏感词库\n10000+ 三级分级词汇\n多模态 NSFW 检测", TEAL),
    ("情绪分析", "S2", "文字 + 语音 + 视频\n多信号情绪融合\n冲突仲裁机制", TEAL_LIGHT),
    ("智能路由", "S3", "规则引擎驱动\n三条干预路线决策\n置信度评分体系", RGBColor(0x1A, 0x8F, 0x9A)),
    ("干预生成", "S4", "危机模板 / 安慰共情\nRAG 知识库检索\nPHQ-9/GAD-7 量表筛查", MINT),
]

card_w = 5.9
card_h = 1.8
positions = [(0.6, 2.1), (6.8, 2.1), (0.6, 4.2), (6.8, 4.2)]

for i, ((title, tag, desc, accent)) in enumerate(cards):
    x, y = positions[i]
    add_rect(slide, x, y, card_w, card_h, CARD_BG, shadow=True)

    # 顶部色条
    add_rect(slide, x, y, card_w, 0.06, accent)

    # 标签圆形
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.3), Inches(y+0.3), Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    set_shape_text(circle, tag, font_size=13, color=WHITE, bold=True)

    add_textbox(slide, x+0.9, y+0.25, 4.5, 0.4, title,
                font_size=19, color=DARK_TEXT, bold=True)
    add_textbox(slide, x+0.9, y+0.7, 4.5, 0.9, desc,
                font_size=12, color=MUTED_TEXT, line_spacing=1.4)

# 技术栈条
tech_y = 6.3
add_textbox(slide, 0.6, tech_y, 0.9, 0.35, "技术栈", font_size=11, color=MUTED_TEXT, bold=True)
techs = [
    ("FastAPI", TEAL), ("Qwen2.5-7B", TEAL_LIGHT), ("SenseVoice", RGBColor(0x1A,0x8F,0x9A)),
    ("Chroma", MINT), ("MySQL", TEAL), ("Edge-TTS", TEAL_LIGHT), ("OpenCV", RGBColor(0x1A,0x8F,0x9A)),
    ("jieba", MINT), ("HSEmotion", TEAL), ("Docker", TEAL_LIGHT)
]
for j, (t, c) in enumerate(techs):
    rx = 1.6 + j * 1.1
    add_rect(slide, rx, tech_y, 1.0, 0.35, c)
    add_textbox(slide, rx, tech_y, 1.0, 0.35, t, font_size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════
# SLIDE 3 — 整体架构图
# ═════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_page_title(slide, 2, "整体系统架构")

# 架构图：前端 → FastAPI → 业务模块 → 基础设施
layers = [
    (0.6, 1.3, 12.1, 0.7, "前端层", "HTML/CSS/JS 聊天界面  ·  语音输入  ·  摄像头拍照  ·  视频上传", TEAL),
    (0.6, 2.5, 12.1, 0.7, "API 网关层", "FastAPI + Uvicorn  ·  CORS  ·  路由分发  ·  请求验证", TEAL_LIGHT),
    (0.6, 3.7, 12.1, 1.4, "业务流水线层", "", RGBColor(0x1A, 0x8F, 0x9A)),
    (0.6, 5.6, 12.1, 0.7, "基础设施层", "MySQL 持久化  ·  Chroma 向量库  ·  Redis 缓存  ·  ffmpeg 预处理", MINT),
]

for (x, y, w, h, label, desc, color) in layers:
    add_rect(slide, x, y, w, h, color)
    add_textbox(slide, x+0.2, y+0.03, 2, 0.35, label, font_size=16, color=WHITE, bold=True)
    if desc:
        add_textbox(slide, x+2.5, y+0.1, w-2.8, 0.5, desc, font_size=12, color=WHITE)
    if label == "业务流水线层":
        # 四个子模块
        sub_modules = ["Safety\n安全过滤", "Emotion\n情绪分析", "Router\n智能路由", "Intervention\n干预生成"]
        for j, sm in enumerate(sub_modules):
            sx = x + 0.3 + j * 3.0
            add_rect(slide, sx, y+0.5, 2.6, 0.75, WHITE)
            add_textbox(slide, sx, y+0.5, 2.6, 0.75, sm, font_size=13, color=color,
                        bold=True, align=PP_ALIGN.CENTER, line_spacing=1.3)
            if j < 3:
                add_arrow_right(slide, sx+2.65, y+0.65, 0.3, 0.22, color)

# 外部服务标注
ext_services = [
    (0.6, 6.7, "LLM 服务 (Qwen2.5-7B)"),
    (3.5, 6.7, "ASR 服务 (SenseVoice/Whisper)"),
    (6.5, 6.7, "TTS 服务 (Edge-TTS)"),
    (9.5, 6.7, "知识库 (Chroma + BM25)"),
]
for (ex, ey, et) in ext_services:
    add_rect(slide, ex, ey, 2.5, 0.4, RGBColor(0xE8, 0xEF, 0xEE))
    add_textbox(slide, ex, ey, 2.5, 0.4, et, font_size=10, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# 箭头连接
add_textbox(slide, 6.2, 2.05, 1.0, 0.4, "▼", font_size=18, color=TEAL, align=PP_ALIGN.CENTER)
add_textbox(slide, 6.2, 3.25, 1.0, 0.4, "▼", font_size=18, color=TEAL_LIGHT, align=PP_ALIGN.CENTER)
add_textbox(slide, 6.2, 5.2, 1.0, 0.4, "▼", font_size=18, color=RGBColor(0x1A,0x8F,0x9A), align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════
# SLIDE 4 — 核心四阶段流水线
# ═════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_page_title(slide, 3, "核心四阶段流水线 (Pipeline)")

# 用户输入标注
add_textbox(slide, 5.6, 1.15, 2.0, 0.4, "用户输入 (文字 / 语音 / 视频)",
            font_size=12, color=MUTED_TEXT, align=PP_ALIGN.CENTER)
add_textbox(slide, 6.3, 1.5, 0.6, 0.3, "▼", font_size=16, color=MUTED_TEXT, align=PP_ALIGN.CENTER)

# 四个阶段 chevron 箭头
stages = [
    ("S1 安全过滤", "关键词匹配\n三级敏感词库\n多模态 NSFW", TEAL),
    ("S2 情绪分析", "文字情绪 LLM\n语音情绪 SenseVoice\n面部情绪 HSEmotion", TEAL_LIGHT),
    ("S3 智能路由", "风险等级评估\n情绪偏向计算\n三路线决策", RGBColor(0x1A, 0x8F, 0x9A)),
    ("S4 干预生成", "危机模板\nLLM 共情对话\nRAG + 量表筛查", MINT),
]

for i, (title, desc, color) in enumerate(stages):
    sx = 0.4 + i * 3.2
    # 阶段框
    add_rect(slide, sx, 2.0, 2.9, 1.8, WHITE, shadow=True)
    add_rect(slide, sx, 2.0, 2.9, 0.06, color)
    add_textbox(slide, sx+0.2, 2.1, 2.5, 0.35, title, font_size=16, color=color, bold=True)
    add_textbox(slide, sx+0.2, 2.5, 2.5, 1.1, desc, font_size=11, color=MUTED_TEXT, line_spacing=1.4)
    if i < 3:
        add_textbox(slide, sx+2.9, 2.55, 0.35, 0.3, "▶", font_size=14, color=MUTED_TEXT, align=PP_ALIGN.CENTER)

# 紧急短路逻辑
add_rect(slide, 0.4, 4.3, 12.5, 2.5, WHITE, shadow=True)
add_rect(slide, 0.4, 4.3, 0.08, 2.5, ACCENT_ORANGE)
add_textbox(slide, 0.7, 4.4, 5, 0.4, "⚠ 紧急短路逻辑", font_size=18, color=ACCENT_ORANGE, bold=True)

# 短路说明
short_text = [
    {"text": "当安全检测发现危险等级 >= 2 时：", "size": 14, "color": DARK_TEXT, "bold": True, "break_line": True},
    {"text": "", "size": 8, "color": MUTED_TEXT, "break_line": True},
    {"text": "1. 跳过阶段二（情绪分析）和阶段三（智能路由）", "size": 13, "color": DARK_TEXT, "break_line": True},
    {"text": "2. 情绪自动标记为 distress（极度痛苦）", "size": 13, "color": DARK_TEXT, "break_line": True},
    {"text": "3. 路由强制设为 crisis（危机干预）", "size": 13, "color": DARK_TEXT, "break_line": True},
    {"text": "4. 直接进入阶段四 → 触发紧急干预模板 + 120 急救推送", "size": 13, "color": DARK_TEXT, "break_line": True},
    {"text": "", "size": 8, "color": MUTED_TEXT, "break_line": True},
    {"text": "设计理念：安全优先。发现危险即跳过所有中间环节，最大限度保障用户生命安全。", "size": 12, "color": ACCENT_ORANGE, "bold": True},
]
add_rich_textbox(slide, 0.7, 4.8, 11.5, 1.8, short_text, line_spacing=1.3)

# 短路箭头可视化
add_rect(slide, 3.6, 3.0, 5.8, 0.04, ACCENT_ORANGE)  # 横线

# ═════════════════════════════════════════════════════════
# SLIDE 5 — 阶段一：安全过滤
# ═════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_page_title(slide, 4, "阶段一：安全过滤 (Safety Filter)")

# 左侧说明
add_textbox(slide, 0.6, 1.2, 5.8, 0.4, "目的：第一时间识别危机内容，触发紧急干预",
            font_size=15, color=DARK_TEXT, bold=True)

detection_items = [
    ("文本关键词匹配", "基于 sensitive_words.json，包含 10000+ 中文敏感词，按 0/1/2 三级分类"),
    ("多模态安全检测", "视频帧检测 · 音频内容审查 · 图像 NSFW 识别"),
    ("模拟 120 急救", "危险等级 ≥ 2 时，自动调用 emergency_push 生成危机响应模板"),
    ("分级阻断策略", "Level 0 放行 · Level 1 标记关注 · Level 2 阻断并触发紧急短路"),
]
for i, (dt, dd) in enumerate(detection_items):
    y = 1.8 + i * 0.9
    add_rect(slide, 0.6, y, 5.8, 0.75, CARD_BG, shadow=True)
    add_textbox(slide, 0.8, y+0.05, 5.4, 0.3, dt, font_size=14, color=TEAL, bold=True)
    add_textbox(slide, 0.8, y+0.38, 5.4, 0.3, dd, font_size=11, color=MUTED_TEXT)

# 右侧统计大卡
add_rect(slide, 7.0, 1.2, 5.8, 3.6, DARK_BG)
add_textbox(slide, 7.3, 1.5, 5.2, 0.4, "敏感词库", font_size=16, color=MINT, bold=True)

# 大数字
stat_y = 2.2
add_textbox(slide, 7.3, stat_y, 2.5, 0.8, "10,000+", font_size=48, color=WHITE, bold=True)
add_textbox(slide, 7.3, stat_y+0.8, 2.5, 0.3, "中文敏感词条", font_size=13, color=MUTED_TEXT)

add_textbox(slide, 10.0, stat_y, 2.5, 0.8, "3 级", font_size=48, color=MINT, bold=True)
add_textbox(slide, 10.0, stat_y+0.8, 2.5, 0.3, "危险等级体系", font_size=13, color=MUTED_TEXT)

# 三级说明
levels = [
    ("Level 0", "一般敏感词，放行", RGBColor(0x5C, 0xD6, 0xC3)),
    ("Level 1", "较重敏感词，标记关注", RGBColor(0xD4, 0xA0, 0x1F)),
    ("Level 2", "极端危险词，阻断+紧急干预", ACCENT_ORANGE),
]
for i, (lv, ld, lc) in enumerate(levels):
    ly = 3.25 + i * 0.4
    add_rect(slide, 7.5, ly, 0.6, 0.3, lc)
    add_textbox(slide, 7.5, ly, 0.6, 0.3, lv, font_size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 8.2, ly, 4.3, 0.3, ld, font_size=12, color=RGBColor(0xBB, 0xCC, 0xD2))

# 输出标注
add_rect(slide, 7.0, 5.1, 5.8, 0.6, TEAL)
add_textbox(slide, 7.2, 5.1, 5.4, 0.6, "输出：SafetyCheckResult { level, blocked, matched_terms }",
            font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# 左下流程
add_rect(slide, 0.6, 5.1, 5.8, 0.6, RGBColor(0xE8, 0xEF, 0xEE))
add_textbox(slide, 0.6, 5.1, 5.8, 0.6, "文本/音频/视频 → 关键词匹配 → 多模态检测 → 等级判定 → 阻断/放行",
            font_size=13, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════
# SLIDE 6 — 阶段二：情绪分析
# ═════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_page_title(slide, 5, "阶段二：情绪分析 (Emotion Analysis)")

# 三列信号源
signals = [
    ("文字情绪", "LLM 情感分析\nEkman 七分类\n情感强度评分", "📝", TEAL),
    ("语音情绪", "SenseVoice 模型\n从音频波形提取\n语调/语速/停顿", "🎤", TEAL_LIGHT),
    ("面部情绪", "HSEmotion 引擎\nMediaPipe 人脸检测\n逐帧表情识别", "📷", RGBColor(0x1A, 0x8F, 0x9A)),
]

for i, (title, desc, icon, color) in enumerate(signals):
    sx = 0.6 + i * 4.2
    add_rect(slide, sx, 1.2, 3.8, 2.2, WHITE, shadow=True)
    add_rect(slide, sx, 1.2, 3.8, 0.06, color)
    # 图标圆形
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(sx+0.3), Inches(1.4), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    set_shape_text(circle, icon, font_size=18, color=WHITE)
    add_textbox(slide, sx+1.0, 1.4, 2.5, 0.4, title, font_size=18, color=DARK_TEXT, bold=True)
    add_textbox(slide, sx+0.3, 2.0, 3.2, 1.2, desc, font_size=12, color=MUTED_TEXT, line_spacing=1.4)

# 汇聚到融合引擎
add_textbox(slide, 2.0, 3.5, 1.0, 0.3, "▼", font_size=20, color=TEAL, align=PP_ALIGN.CENTER)
add_textbox(slide, 6.0, 3.5, 1.0, 0.3, "▼", font_size=20, color=TEAL_LIGHT, align=PP_ALIGN.CENTER)
add_textbox(slide, 10.0, 3.5, 1.0, 0.3, "▼", font_size=20, color=RGBColor(0x1A,0x8F,0x9A), align=PP_ALIGN.CENTER)

# 融合引擎
add_rect(slide, 3.5, 3.9, 6.3, 1.1, DARK_BG)
add_textbox(slide, 3.7, 3.95, 5.9, 0.35, "Emotion Fusion 情绪融合引擎",
            font_size=18, color=MINT, bold=True)
add_textbox(slide, 3.7, 4.35, 5.9, 0.55,
            "多信号加权综合  ·  冲突仲裁（按置信度裁决）  ·  情绪标签映射（Ekman → 8 类合约标签）",
            font_size=12, color=RGBColor(0xBB, 0xCC, 0xD2))

# 输出
add_textbox(slide, 6.0, 5.1, 1.3, 0.3, "▼", font_size=18, color=MINT, align=PP_ALIGN.CENTER)
add_rect(slide, 3.0, 5.5, 7.3, 0.7, TEAL)
add_textbox(slide, 3.2, 5.5, 6.9, 0.7,
            "输出：EmotionTags { primary_emotion, intensity, risk, modality_notes }",
            font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Ekman 情绪标注
emotions = ["愤怒", "厌恶", "恐惧", "快乐", "悲伤", "惊讶", "中性"]
for i, e in enumerate(emotions):
    ex = 1.2 + i * 1.75
    add_rect(slide, ex, 6.5, 1.5, 0.4, RGBColor(0xE8, 0xEF, 0xEE))
    add_textbox(slide, ex, 6.5, 1.5, 0.4, e, font_size=10, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════
# SLIDE 7 — 阶段三：智能路由
# ═════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_page_title(slide, 6, "阶段三：智能路由 (Route Decision)")

add_textbox(slide, 0.6, 1.2, 8, 0.4, "目的：综合安全评估与情绪分析结果，决定干预路线",
            font_size=15, color=DARK_TEXT, bold=True)

# 决策表
table_data = [
    ["条件", "路由结果", "说明", "置信度"],
    ["安全等级 >= 2", "CRISIS 危机", "直接升级，最高优先级", "强制 1.0"],
    ["情绪强度 >= 0.8 + 高风险", "CRISIS 危机", "高情绪风险触发危机", "0.8 - 1.0"],
    ["情绪强度 0.4~0.8 / 负面", "COMFORT 安慰", "需要共情安抚", "0.6 - 0.9"],
    ["情绪强度 < 0.4 / 正面", "KNOWLEDGE 知识", "理性交流，知识科普", "0.5 - 0.9"],
    ["用户主动询问知识", "KNOWLEDGE 知识", "满足求知需求", "0.7 - 1.0"],
]

table = slide.shapes.add_table(
    len(table_data), 4,
    Inches(0.6), Inches(1.8),
    Inches(12.0), Inches(2.0)
).table

# 表格样式
header_colors = [TEAL] * 4
body_colors = [
    [RGBColor(0xFF, 0xF0, 0xEC), RGBColor(0xFF, 0xF0, 0xEC), ACCENT_ORANGE, RGBColor(0xFF, 0xF0, 0xEC)],
    [RGBColor(0xFF, 0xF5, 0xF0), RGBColor(0xFF, 0xF5, 0xF0), ACCENT_ORANGE, RGBColor(0xFF, 0xF5, 0xF0)],
    [WHITE, WHITE, ROUTE_COMFORT, WHITE],
    [RGBColor(0xF0, 0xF8, 0xF9), RGBColor(0xF0, 0xF8, 0xF9), ROUTE_KNOWLEDGE, RGBColor(0xF0, 0xF8, 0xF9)],
    [WHITE, WHITE, ROUTE_KNOWLEDGE, WHITE],
]

for i, row_data in enumerate(table_data):
    for j, cell_text in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.name = "Microsoft YaHei"
            paragraph.alignment = PP_ALIGN.CENTER
            if i == 0:
                paragraph.font.color.rgb = WHITE
                paragraph.font.bold = True
            else:
                if j == 2:
                    paragraph.font.color.rgb = body_colors[i-1][j]
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = DARK_TEXT
        if body_colors[i-1][j] if i > 0 else None:
            if i > 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = body_colors[i-1][j]
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TEAL

# 置信度计算公式
add_rect(slide, 0.6, 4.2, 12.0, 0.7, WHITE, shadow=True)
add_textbox(slide, 0.8, 4.25, 11.6, 0.6,
            "置信度计算：基础分 - 边界模糊惩罚 - 信号冲突惩罚 = 最终置信度",
            font_size=14, color=DARK_TEXT, bold=True)

# 三条路线可视化
routes = [
    ("COMFORT 安慰路线", "以 LLM 共情对话为主\n提供心理安抚与情绪疏导", ROUTE_COMFORT),
    ("KNOWLEDGE 知识路线", "RAG 知识库检索 + 量表筛查\n提供专业心理知识", ROUTE_KNOWLEDGE),
    ("CRISIS 危机路线", "确定性模板回复 + 紧急推送\n保障用户生命安全", ROUTE_CRISIS),
]

for i, (rt, rd, rc) in enumerate(routes):
    rx = 0.6 + i * 4.2
    add_rect(slide, rx, 5.2, 3.9, 1.5, WHITE, shadow=True)
    add_rect(slide, rx, 5.2, 0.08, 1.5, rc)
    add_textbox(slide, rx+0.3, 5.35, 3.3, 0.35, rt, font_size=15, color=rc, bold=True)
    add_textbox(slide, rx+0.3, 5.75, 3.3, 0.8, rd, font_size=11, color=MUTED_TEXT, line_spacing=1.4)

# ═════════════════════════════════════════════════════════
# SLIDE 8 — 阶段四：干预生成
# ═════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_page_title(slide, 7, "阶段四：干预生成 (Intervention)")

# 三条干预路线大卡片
branches = [
    ("CRISIS 危机干预", ROUTE_CRISIS, [
        "匹配危机类型模板（自杀/自伤/暴力/急症）",
        "生成确定性回复（不经过 LLM，保证稳定可靠）",
        "触发紧急推送（模拟 120 救援 API）",
        "输出：安抚话术 + 求助信息 + 紧急联系人",
    ]),
    ("COMFORT 安慰干预", ROUTE_COMFORT, [
        "调用 LLM（Qwen2.5-7B）生成共情回复",
        "使用安抚提示词模板 + 8 种心理治疗技术随机选择",
        "倾听与共情 · 认知重构 · 行为激活 · 放松训练",
        "接纳承诺 · 积极心理 · 自我关怀 · 社交支持",
    ]),
    ("KNOWLEDGE 知识干预", ROUTE_KNOWLEDGE, [
        "量表筛查判断：PHQ-9 抑郁筛查 / GAD-7 焦虑筛查",
        "RAG 检索：LLM 查询分类 → Chroma 密集向量 + BM25 关键词",
        "RRF 融合排序 → Top 3 → LLM 结合检索结果生成回复",
        "支持多轮量表流程：开始 → 逐题施测 → 评分反馈",
    ]),
]

for i, (title, color, items) in enumerate(branches):
    bx = 0.45 + i * 4.25
    add_rect(slide, bx, 1.2, 4.05, 4.8, WHITE, shadow=True)
    add_rect(slide, bx, 1.2, 4.05, 0.6, color)
    add_textbox(slide, bx+0.2, 1.25, 3.65, 0.5, title, font_size=17, color=WHITE, bold=True)

    for j, item in enumerate(items):
        iy = 2.0 + j * 0.75
        # 小圆点
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(bx+0.2), Inches(iy+0.08), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        add_textbox(slide, bx+0.45, iy, 3.4, 0.6, item, font_size=11, color=DARK_TEXT, line_spacing=1.3)

# 下方小结
add_rect(slide, 0.6, 6.3, 12.1, 0.5, DARK_BG)
add_textbox(slide, 0.8, 6.3, 11.7, 0.5,
            "最终输出：InterventionResult { reply, empathy, suggestion, action_items } → 返回给用户",
            font_size=14, color=MINT, bold=True, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════
# SLIDE 9 — 多模态处理流程
# ═════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_page_title(slide, 8, "多模态处理流程：视频输入")

# 流程图横排
flow_steps = [
    ("视频文件", "上传 .mp4", TEAL),
    ("ffmpeg\n预处理", "提取音频\n逐帧采样", TEAL_LIGHT),
    ("ASR\n语音识别", "SenseVoice\n→ 文字", RGBColor(0x1A, 0x8F, 0x9A)),
    ("音频情绪", "SenseVoice\n语音情绪", RGBColor(0x1A, 0x8F, 0x9A)),
    ("面部情绪", "HSEmotion\n+ MediaPipe", RGBColor(0x1A, 0x8F, 0x9A)),
    ("情绪融合", "加权综合\n冲突仲裁", MINT),
    ("四阶段\n流水线", "Safety → ...\n→ Intervention", TEAL),
]

for i, (title, desc, color) in enumerate(flow_steps):
    fx = 0.35 + i * 1.85
    fw = 1.6 if i < 6 else 1.72
    add_rect(slide, fx, 1.5, fw, 1.3, color)
    add_textbox(slide, fx, 1.6, fw, 0.4, title, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, fx, 2.05, fw, 0.6, desc, font_size=10, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 6:
        add_textbox(slide, fx+fw, 1.85, 0.3, 0.3, "▶", font_size=11, color=MUTED_TEXT, align=PP_ALIGN.CENTER)

# 详细说明区
details = [
    ("VideoPreprocessor 模块", [
        "使用 ffmpeg 从视频文件中提取 WAV 音频",
        "使用 OpenCV (cv2.VideoCapture) 按时间均匀采帧（最多 20 帧）",
        "MediaPipe 人脸检测 → 逐帧裁剪人脸区域",
        "最少 3 个有效帧 + 人脸检出率 ≥ 30% 才进入情绪分析",
    ]),
    ("情绪融合机制 (EmotionFusion)", [
        "文字情绪 + 语音情绪 + 面部情绪 → 三路信号加权融合",
        "各通道置信度不同时，按置信度加权取主情绪",
        "若某通道信号缺失或不可靠（如无人脸），自动剔除并降权",
        "最终输出统一的 EmotionTags 给下游路由模块",
    ]),
]

for j, (dh, ditems) in enumerate(details):
    dx = 0.6 + j * 6.2
    add_rect(slide, dx, 3.2, 6.0, 3.0, WHITE, shadow=True)
    add_rect(slide, dx, 3.2, 6.0, 0.06, TEAL if j == 0 else MINT)
    add_textbox(slide, dx+0.2, 3.3, 5.6, 0.35, dh, font_size=16, color=TEAL, bold=True)
    for k, di in enumerate(ditems):
        add_textbox(slide, dx+0.4, 3.8+k*0.5, 5.2, 0.4,
                    f"• {di}", font_size=12, color=DARK_TEXT, line_spacing=1.3)

# ═════════════════════════════════════════════════════════
# SLIDE 10 — API 路由结构
# ═════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_page_title(slide, 9, "API 路由结构")

api_data = [
    ["路由", "方法", "功能", "说明"],
    ["/api/v1/pipeline/run", "POST", "文本流水线", "输入文字 → 完整四阶段处理"],
    ["/api/v1/pipeline/run-with-video", "POST", "视频流水线", "上传视频 → 预处理 → 四阶段"],
    ["/api/v1/multimodal/asr", "POST", "语音转文字", "独立 ASR 服务 (SenseVoice/Whisper)"],
    ["/api/v1/multimodal/tts", "POST", "文字转语音", "Edge-TTS 语音合成"],
    ["/api/v1/multimodal/emotion-detect", "POST", "情绪检测", "图片/音频情绪独立分析"],
    ["/api/v1/multimodal/safety-check", "POST", "安全检查", "独立安全过滤接口"],
    ["/api/v1/modules/safety", "POST", "安全模块", "独立模块 API，支持 Mock 切换"],
    ["/api/v1/modules/emotion", "POST", "情绪模块", "独立模块 API，支持 Mock 切换"],
    ["/api/v1/modules/router", "POST", "路由模块", "独立模块 API，支持 Mock 切换"],
    ["/api/v1/modules/intervention", "POST", "干预模块", "独立模块 API，支持 Mock 切换"],
    ["/api/v1/chat", "POST", "旧版聊天", "已标记废弃，保留兼容"],
    ["/api/v1/chat/stream", "POST", "流式聊天", "SSE 实时流式输出"],
]

table = slide.shapes.add_table(
    len(api_data), 4,
    Inches(0.6), Inches(1.2),
    Inches(12.0), Inches(5.0)
).table

col_widths = [4.5, 1.2, 2.0, 4.3]
for j, w in enumerate(col_widths):
    table.columns[j].width = Inches(w)

row_colors = [WHITE, RGBColor(0xF5, 0xFA, 0xF9)]
for i, row_data in enumerate(api_data):
    for j, cell_text in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.name = "Microsoft YaHei"
            if i == 0:
                paragraph.font.color.rgb = WHITE
                paragraph.font.bold = True
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                paragraph.font.color.rgb = DARK_TEXT
                if j == 0:
                    paragraph.font.name = "Consolas"
                    paragraph.font.size = Pt(10)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        # 填充色
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BG
        elif i % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF5, 0xFA, 0xF9)

# Mock 说明
add_rect(slide, 0.6, 6.4, 12.0, 0.5, RGBColor(0xE8, 0xEF, 0xEE))
add_textbox(slide, 0.8, 6.4, 11.6, 0.5,
            "各模块通过 .env 的 MOCK_* 开关支持独立 Mock/Real 模式切换，实现并行开发与独立测试",
            font_size=12, color=MUTED_TEXT, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════
# SLIDE 11 — 结束页
# ═════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_rect(slide, 0, 0, 13.333, 0.06, MINT)

add_textbox(slide, 1, 2.0, 11.3, 1.0,
            "感谢聆听",
            font_size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_rect(slide, 5.5, 3.1, 2.3, 0.04, MINT)

add_textbox(slide, 1, 3.5, 11.3, 0.6,
            "心理异常智能早筛与精准干预系统",
            font_size=22, color=MINT, align=PP_ALIGN.CENTER)

add_textbox(slide, 1, 4.3, 11.3, 0.5,
            "安全过滤  ·  情绪分析  ·  智能路由  ·  精准干预",
            font_size=16, color=RGBColor(0x9B, 0xB8, 0xBF), align=PP_ALIGN.CENTER)

add_textbox(slide, 1, 5.5, 11.3, 0.5,
            "Q & A",
            font_size=20, color=RGBColor(0x7A, 0x9E, 0xA8), align=PP_ALIGN.CENTER)

# ── 保存 ──────────────────────────────────────────────
output_path = r"c:\Users\xt\Desktop\Project\mental-intervene-master\心理异常智能早筛与精准干预系统.pptx"
prs.save(output_path)
print(f"PPT 已生成: {output_path}")
print(f"共 {len(prs.slides)} 页幻灯片")
